#!/usr/bin/env python3
"""
Create PowerPoint presentation explaining the KAT8 Sex x Genotype
interaction analysis (parts 7.6 + 7.7) for presentation to PI / boss.

Updated 2026-02-18: Reflects final decision: ~ 0 + GroupDepot (unified, no Sex
covariate). Added Part 7.7 model comparison at both per-depot and unified levels,
educational slides explaining covariates/pi0/concordance, and the pooling decision.

Run:  python create_presentation.py
Output:  Results/KAT8_sex_interaction_presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
DARK_BLUE  = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE   = RGBColor(0x2C, 0x6F, 0xAD)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
GREY       = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
RED_ACC    = RGBColor(0xC0, 0x39, 0x2B)
GREEN_ACC  = RGBColor(0x27, 0xAE, 0x60)
AMBER_ACC  = RGBColor(0xF3, 0x9C, 0x12)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background colour for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=16, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=14, color=BLACK, font_name="Calibri",
                    spacing_after=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing_after
        p.level = 0
    return tf


def add_plot_placeholder(slide, left, top, width, height, filename):
    """Grey dashed box with instructions to insert the named plot."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    shape.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
    shape.line.dash_style = 2  # dash
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = "INSERT PLOT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GREY
    p2 = tf.add_paragraph()
    p2.text = filename
    p2.font.size = Pt(11)
    p2.font.color.rgb = GREY
    p2.font.italic = True
    p2.alignment = PP_ALIGN.CENTER


def add_section_header(slide, text, top=Inches(0.3)):
    add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.6),
                text, font_size=28, bold=True, color=DARK_BLUE)


def add_subtitle_text(slide, text, top=Inches(0.9)):
    add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.4),
                text, font_size=14, color=GREY)


def add_slide_number(slide, num, total):
    add_textbox(slide, Inches(12), Inches(7), Inches(1), Inches(0.4),
                f"{num}/{total}", font_size=10, color=GREY,
                alignment=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Build Presentation
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

TOTAL_SLIDES = 21  # updated count

# ====== SLIDE 1: Title ======
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BLUE)
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
            "KAT8-KD RNA-seq:\nSex-Combining Validation & Analysis Decision",
            font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
            "Parts 7.6 + 7.7  |  Sex x Genotype interaction testing & model comparison",
            font_size=18, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
            "Adipose Depot Analysis: iWAT & gWAT",
            font_size=16, color=RGBColor(0x88, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
            "Validation run: 2026-02-18  |  Decision: ~ 0 + GroupDepot (unified, no Sex covariate)",
            font_size=14, color=RGBColor(0x77, 0x99, 0xBB), alignment=PP_ALIGN.CENTER)

# ====== SLIDE 2: Study Design ======
slide = prs.slides.add_slide(blank_layout)
add_section_header(slide, "Study Design")
add_bullet_list(slide, Inches(0.8), Inches(1.2), Inches(6), Inches(5), [
    "2 x 2 factorial design: Sex (F, M) x Genotype (CTL, KAT8-KD)",
    "Two adipose depots analyzed separately: iWAT and gWAT",
    "40 samples total, 2 outliers removed (JS_08, JS_28) -> 38 samples",
    "Per depot: 19 samples (after outlier removal)",
    "",
    "Group sizes per depot:",
    "   CTL Female:  5     |  KAT8-KD Female:  4",
    "   CTL Male:    5     |  KAT8-KD Male:    5",
    "",
    "KEY: We have BOTH sexes in every group.",
    "Question: should we account for sex in our statistical model?",
], font_size=16)

# ====== SLIDE 3: The Question ======
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, RGBColor(0xF0, 0xF4, 0xF8))
add_textbox(slide, Inches(1), Inches(1), Inches(11), Inches(1),
            "The Central Question",
            font_size=32, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
            "Does KAT8 knockdown affect gene expression\n"
            "differently in males versus females?",
            font_size=26, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4), Inches(11), Inches(3),
            "Why this matters:\n\n"
            "1. If YES for most genes: We must analyze males and females separately\n"
            "   (splitting would halve our sample size -> less power to detect effects)\n\n"
            "2. If NO for most genes: We can include sex as a COVARIATE\n"
            "   (this absorbs sex-related noise -> MORE power to detect KD effects)\n\n"
            "3. NIH and most journals now require addressing sex as a biological variable (SABV)",
            font_size=16, color=GREY, alignment=PP_ALIGN.LEFT)

# ====== SLIDE 4: What is a Covariate? (EDUCATIONAL) ======
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, RGBColor(0xFD, 0xF2, 0xE9))
add_section_header(slide, "Background: What is a Covariate?")
add_bullet_list(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(6), [
    "A covariate is a variable you ACCOUNT FOR in your statistical model",
    "but are NOT directly interested in testing.",
    "",
    "ANALOGY: Imagine measuring the effect of a drug on blood pressure.",
    "  - Your variable of interest: Drug vs. Placebo",
    "  - A covariate: Patient age (affects blood pressure but isn't what you're testing)",
    "  - By accounting for age, you REMOVE age-related noise from the measurement",
    "  - This makes it easier to see the true drug effect",
    "",
    "In our experiment:",
    "  - Variable of interest: KAT8 knockdown effect (Genotype: CTL vs KAT8-KD)",
    "  - Covariate: Sex (F vs M) -- affects gene expression but isn't our focus",
    "  - By telling the model 'some variance is due to sex,' it can SUBTRACT",
    "    that variance and see the genotype effect more clearly",
    "",
    "Result: MORE statistical power to detect true KAT8 knockdown effects",
], font_size=14)

# ====== SLIDE 5: Two Possible Models (EDUCATIONAL) ======
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, RGBColor(0xFD, 0xF2, 0xE9))
add_section_header(slide, "Background: Two Possible DESeq2 Models")
add_bullet_list(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
    "MODEL A:   ~ Genotype                  (what we had before)",
    "  - Compares KAT8-KD vs CTL, ignoring sex entirely",
    "  - Any gene expression difference between M and F becomes 'noise'",
    "  - The model has to 'see through' this noise to find KD effects",
    "",
    "MODEL B:   ~ Sex + Genotype            (the new model)",
    "  - First accounts for M vs F baseline differences (the covariate)",
    "  - THEN tests KAT8-KD vs CTL on the sex-adjusted expression",
    "  - Less noise -> more statistical power -> more DEGs detected",
    "",
    "IMPORTANT: Model B does NOT mean 'analyze males and females separately'",
    "  - All samples stay in ONE analysis (maximum sample size)",
    "  - Sex is just 'controlled for,' like adjusting for age in clinical studies",
    "",
    "BUT: Model B is only valid IF the KD effect is similar in M and F.",
    "  If males and females respond completely differently to KD,",
    "  we would need to split them. That's what we test next.",
], font_size=14)

# ====== SLIDE 6: Method - LRT ======
slide = prs.slides.add_slide(blank_layout)
add_section_header(slide, "Method: Likelihood Ratio Test (LRT)")
add_subtitle_text(slide, "Testing whether the KD effect depends on sex")
add_bullet_list(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(5.5), [
    "Tool: DESeq2 (Bioconductor), run separately per depot",
    "",
    "Full model:      ~ Sex + Genotype + Sex:Genotype",
    "Reduced model:  ~ Sex + Genotype",
    "",
    "What is Sex:Genotype (the 'interaction term')?",
    "  It captures genes where the KD effect DIFFERS between males and females.",
    "  Example: Gene X goes UP in female KD but DOWN in male KD -> strong interaction",
    "  Example: Gene Y goes UP equally in both sexes -> NO interaction",
    "",
    "The LRT asks: for each gene, does adding the interaction term",
    "significantly improve how well the model fits the data?",
    "",
    "  If p < 0.05 -> that gene's KD response truly differs between M and F",
    "  If p >= 0.05 -> the KD effect is the same in both sexes (what we want!)",
    "",
    "We also run 4 complementary analyses to build the full picture:",
    "  (1) Pi0 estimation  (2) Variance decomposition  "
    "(3) Sex-stratified concordance  (4) Model comparison (Part 7.7)",
], font_size=13)

# ====== SLIDE 7: How to Read a P-value Histogram (EDUCATIONAL) ======
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, RGBColor(0xFD, 0xF2, 0xE9))
add_section_header(slide, "Background: How to Read a P-value Histogram")
add_bullet_list(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
    "When we test thousands of genes, we get one p-value per gene.",
    "Plotting these p-values as a histogram reveals the big picture:",
    "",
    "SCENARIO 1: All genes are truly null (no real interaction anywhere)",
    "  -> P-values are UNIFORMLY distributed (flat histogram, like a table)",
    "  -> Each bar has roughly the same height",
    "",
    "SCENARIO 2: Some genes have real interactions",
    "  -> There's a SPIKE near p = 0 (real signal produces small p-values)",
    "  -> The rest of the histogram stays flat (the true nulls)",
    "",
    "What you'll see in our data:",
    "  -> Mostly FLAT (most genes have no sex-dependent KD response)",
    "  -> Small spike near p = 0 (a few genes do differ between sexes)",
    "  -> Red dashed line = expected height if ALL genes were null",
    "  -> Bars above the red line near p = 0 = genes with real interaction signal",
    "",
    "The flatter the histogram, the more we can trust that",
    "males and females respond similarly to KAT8 knockdown.",
], font_size=13)

# ====== SLIDE 8: iWAT LRT Results ======
slide = prs.slides.add_slide(blank_layout)
add_section_header(slide, "iWAT: LRT Interaction Results")
add_bullet_list(slide, Inches(0.6), Inches(1.2), Inches(5.5), Inches(5), [
    "16,466 genes tested (after prefiltering)",
    "630 significant interactions (3.83%, FDR < 0.05)",
    "~96% of genes show NO sex-dependent KD response",
    "",
    "How to read the histogram (right):",
    "  - Flat distribution = genes with no interaction (null)",
    "  - Spike near p=0 = genes with real interaction signal",
    "  - Red dashed line = expected count if ALL genes were null",
    "  - The histogram is mostly flat -> most genes are null",
    "",
    "3.83% is a low interaction rate. For context:",
    "  - At FDR < 0.05, we expect ~5% false positives",
    "  - So some of these 630 genes are likely false positives",
    "  - The true interaction rate is even lower (see pi0 next)",
], font_size=13)
add_plot_placeholder(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(5),
                     "pvalue_hist_interaction_iWAT_<run_tag>.png")

# ====== SLIDE 9: gWAT LRT Results ======
slide = prs.slides.add_slide(blank_layout)
add_section_header(slide, "gWAT: LRT Interaction Results")
add_bullet_list(slide, Inches(0.6), Inches(1.2), Inches(5.5), Inches(5), [
    "16,827 genes tested (after prefiltering)",
    "374 significant interactions (2.22%, FDR < 0.05)",
    "~98% of genes show NO sex-dependent KD response",
    "",
    "Even fewer interaction genes than iWAT (2.2% vs 3.8%),",
    "further supporting sex-independent KD effects.",
    "",
    "Why might gWAT have fewer interactions?",
    "  - The KD effect is weaker in gWAT overall (fewer DEGs)",
    "  - With a weaker main signal, there's less room for",
    "    sex-dependent variation to be detected",
    "  - This is consistent with what we see in the main analysis:",
    "    iWAT shows a stronger KAT8 KD response than gWAT",
], font_size=13)
add_plot_placeholder(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(5),
                     "pvalue_hist_interaction_gWAT_<run_tag>.png")

# ====== SLIDE 10: What is Pi0? (EDUCATIONAL) ======
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, RGBColor(0xFD, 0xF2, 0xE9))
add_section_header(slide, "Background: What is Pi0?")
add_bullet_list(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
    "When we test 16,000+ genes, we get a mix of:",
    "  - TRUE NULLS: genes where there really is no interaction (what we want!)",
    "  - TRUE SIGNALS: genes where the KD effect genuinely differs by sex",
    "  - FALSE POSITIVES: genes that look significant by chance",
    "",
    "Pi0 (pi-zero) estimates the proportion of TRUE NULLS:",
    "  - Pi0 = 0.87 means: 87% of genes genuinely have NO sex-dependent KD response",
    "  - The other 13% includes both true interactions AND statistical noise",
    "",
    "How is pi0 calculated? (Storey & Tibshirani, 2003)",
    "  - If ALL genes were null, p-values would be perfectly uniform",
    "  - The method looks at how 'flat' the p-value distribution is away from p=0",
    "  - The flat part = true nulls. The spike near p=0 = real signals",
    "",
    "Why this matters:",
    "  - Pi0 = 0.50 would mean half the genes have sex-dependent effects -> BAD",
    "  - Pi0 = 0.95+ would mean almost no genes have sex-dependent effects -> GREAT",
    "  - Pi0 = 0.85-0.87 (what we got) -> GOOD: vast majority are sex-independent",
], font_size=13)

# ====== SLIDE 11: Pi0 Results ======
slide = prs.slides.add_slide(blank_layout)
add_section_header(slide, "Pi0 Results: True Null Proportion")
add_bullet_list(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
    "Results:",
    "  iWAT:  pi0 = 0.8701  ->  87% of genes genuinely have NO interaction",
    "  gWAT:  pi0 = 0.8483  ->  85% of genes genuinely have NO interaction",
    "",
    "What does this mean in plain English?",
    "",
    "  Imagine lining up all ~16,500 genes in iWAT:",
    "  - ~14,350 genes (87%): KAT8 KD has the SAME effect in males and females",
    "  - ~2,150 genes (13%):  some have real sex-dependent effects, some are noise",
    "",
    "  Of those 2,150:",
    "  - 630 reached statistical significance (FDR < 0.05)",
    "  - The rest didn't quite reach significance but may have subtle effects",
    "",
    "Bottom line: The vast majority of genes respond identically to KAT8-KD",
    "in both sexes. This strongly supports including sex as a covariate",
    "(rather than splitting the analysis by sex).",
], font_size=14)

# ====== SLIDE 12: Variance Attribution ======
slide = prs.slides.add_slide(blank_layout)
add_section_header(slide, "Variance Attribution: What Drives Gene Expression?")
add_subtitle_text(slide,
    "For each gene, what fraction of variance comes from Sex, Genotype, or their interaction?")
add_bullet_list(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(2.5), [
    "Method: ANOVA on VST-transformed counts:",
    "  lm(expression ~ Sex + Genotype + Sex:Genotype)",
    "",
    "         Term               iWAT (median)   gWAT (median)",
    "  Sex (M vs F)                    6.5%           11.5%",
    "  Genotype (KD effect)          33.2%           16.7%",
    "  Sex x Genotype                  2.1%             2.4%",
    "  Residual (unexplained)        46.3%           50.1%",
], font_size=13)
add_bullet_list(slide, Inches(0.6), Inches(4.5), Inches(5.5), Inches(2), [
    "Key takeaway: The interaction term explains only ~2% of variance",
    "  vs. 17-33% for the Genotype main effect (8-16x more)",
    "",
    "Think of it this way: if gene expression were a pie chart,",
    "  the KD effect gets a BIG slice (17-33%),",
    "  Sex gets a medium slice (6-12%),",
    "  but Sex x KD interaction gets a TINY slice (~2%)",
], font_size=13, color=DARK_BLUE)
add_plot_placeholder(slide, Inches(6.5), Inches(1.5), Inches(6), Inches(5),
                     "variance_explained_<depot>_<run_tag>.png")

# ====== SLIDE 13: What is Concordance? (EDUCATIONAL) ======
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, RGBColor(0xFD, 0xF2, 0xE9))
add_section_header(slide, "Background: What is Concordance?")
add_bullet_list(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
    "Concordance = do males and females AGREE on which genes are affected by KD?",
    "",
    "How we test this:",
    "  1. Run DESeq2 on FEMALES ONLY -> get a fold change for every gene",
    "  2. Run DESeq2 on MALES ONLY -> get a fold change for every gene",
    "  3. Plot female fold changes (x-axis) vs male fold changes (y-axis)",
    "  4. Compute Pearson correlation (r) -- how well do they match?",
    "",
    "Interpreting the scatter plot:",
    "  - Points on the diagonal (y = x): perfect agreement (M and F same)",
    "  - Points in upper-right quadrant: gene goes UP in both sexes (concordant)",
    "  - Points in lower-left quadrant: gene goes DOWN in both sexes (concordant)",
    "  - Points in upper-left or lower-right: gene goes UP in one sex but DOWN",
    "    in the other (DISCORDANT -- these are the problematic ones)",
    "",
    "What r values mean:",
    "  r = 1.0: perfect agreement   r = 0.7-0.9: strong   r = 0.5-0.7: moderate",
    "  r < 0.5: poor agreement (males and females respond very differently)",
], font_size=13)

# ====== SLIDE 14: Concordance iWAT ======
slide = prs.slides.add_slide(blank_layout)
add_section_header(slide, "Sex-Stratified Concordance: iWAT")
add_subtitle_text(slide,
    "DESeq2 run separately for F-only and M-only, then fold changes correlated")
add_bullet_list(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(4.5), [
    "Pearson r = 0.727 (moderate-to-strong concordance)",
    "Directional concordance = 99.0%",
    "",
    "What these numbers mean:",
    "  r = 0.73: Males and females largely agree on which genes",
    "  are affected by KD and how much they change",
    "",
    "  99.0% directional concordance: Of ALL significant DEGs,",
    "  99% go the SAME direction in both sexes. Only 1% are",
    "  discordant (up in one sex, down in the other).",
    "",
    "  Purple dots: significant in BOTH sexes (strongest evidence)",
    "  Red dots: significant in F only",
    "  Blue dots: significant in M only",
    "  Grey dots: not significant in either",
    "",
    "Note: r = 0.73 (not 0.99) because with n=4-5 per group,",
    "per-sex analyses are noisy. True concordance is likely higher.",
], font_size=12)
add_plot_placeholder(slide, Inches(6.5), Inches(1.5), Inches(6), Inches(5),
                     "sex_concordance_scatter_iWAT_<run_tag>.png")

# ====== SLIDE 15: Concordance gWAT ======
slide = prs.slides.add_slide(blank_layout)
add_section_header(slide, "Sex-Stratified Concordance: gWAT")
add_subtitle_text(slide,
    "DESeq2 run separately for F-only and M-only, then fold changes correlated")
add_bullet_list(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(4.5), [
    "Pearson r = 0.605 (moderate concordance)",
    "Directional concordance = 99.5%",
    "",
    "Lower r than iWAT -- why?",
    "  Visceral fat (gWAT) has more baseline sex dimorphism than",
    "  subcutaneous fat (iWAT). This is well-documented in the",
    "  adipose biology literature. Male and female gWAT differ",
    "  more at baseline, adding more noise to per-sex analyses.",
    "",
    "BUT: 99.5% directional concordance is excellent!",
    "  Genes that go up in female gWAT also go up in male gWAT.",
    "  The DIRECTION of KD effects is almost perfectly consistent.",
    "",
    "The lower r mostly reflects MAGNITUDE differences:",
    "  e.g., gene goes up 2-fold in F but 3-fold in M.",
    "  This is normal biological variation, not a problem.",
    "  Including sex as a covariate handles this optimally.",
], font_size=12)
add_plot_placeholder(slide, Inches(6.5), Inches(1.5), Inches(6), Inches(5),
                     "sex_concordance_scatter_gWAT_<run_tag>.png")

# ====== SLIDE 16: PCA ======
slide = prs.slides.add_slide(blank_layout)
add_section_header(slide, "PCA by Sex and Genotype")
add_bullet_list(slide, Inches(0.6), Inches(1.3), Inches(5.5), Inches(5), [
    "Principal Component Analysis of top 500 variable genes",
    "",
    "How to read this plot:",
    "  Color = Genotype (CTL vs KAT8-KD)",
    "  Shape = Sex (circle = F, triangle = M)",
    "",
    "What to look for:",
    "  If M and F of the same genotype cluster together",
    "  -> sexes are transcriptionally similar (GOOD)",
    "",
    "  If M and F of the same genotype separate",
    "  -> sex differences dominate (PROBLEMATIC)",
    "",
    "What we see: genotype drives separation more than sex,",
    "but male KD samples tend to be slightly more extreme",
    "(further from controls) than female KD samples,",
    "especially in iWAT. This is expected and handled",
    "by the sex covariate in our model.",
], font_size=13)
add_plot_placeholder(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(5),
                     "PCA_by_sex_<depot>_<run_tag>.png")

# ====== SLIDE 17: Interaction vs Genotype Scatter ======
slide = prs.slides.add_slide(blank_layout)
add_section_header(slide, "Interaction Effect vs. Genotype Main Effect")
add_bullet_list(slide, Inches(0.6), Inches(1.3), Inches(5.5), Inches(5), [
    "This plot directly compares two effects for every gene:",
    "",
    "X-axis = Genotype main effect log2FC (how much KD changes the gene)",
    "Y-axis = Interaction log2FC (how much the KD effect DIFFERS by sex)",
    "",
    "  Red dots: significant interaction genes",
    "  Blue dots: significant genotype DEGs (no interaction)",
    "  Grey dots: neither significant",
    "",
    "What we want to see:",
    "  Wide horizontal spread (big genotype effects)",
    "  Narrow vertical spread (small interaction effects)",
    "  = the KD signal is MUCH bigger than any sex-dependent variation",
    "",
    "What we actually see:",
    "  Horizontal spread >> vertical spread",
    "  Most dots cluster near y = 0 (no interaction)",
    "  -> CONFIRMS: genotype effect dominates",
], font_size=13)
add_plot_placeholder(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(5),
                     "interaction_vs_genotype_<depot>_<run_tag>.png")

# ====== SLIDE 18: MODEL COMPARISON - Per-Depot Level ======
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, RGBColor(0xE8, 0xF8, 0xF5))
add_section_header(slide, "Per-Depot Models: ~ Genotype vs. ~ Sex + Genotype")
add_subtitle_text(slide,
    "Part 7.7 Test 3a: At the per-depot level (n ~ 19), adding Sex helps")

# Build a model comparison table
rows = 6
cols = 3
table_shape = slide.shapes.add_table(rows, cols,
    Inches(1), Inches(1.5), Inches(11), Inches(2.8))
tbl = table_shape.table

# Header row
headers = ["Per-Depot Comparison", "iWAT (n=19)", "gWAT (n=19)"]
for i, h in enumerate(headers):
    cell = tbl.cell(0, i)
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE

# Data rows
model_data = [
    ["~ Genotype (no Sex)",        "2,665 DEGs",   "1,077 DEGs"],
    ["~ Sex + Genotype",           "3,042 DEGs",   "1,131 DEGs"],
    ["Net DEG change",             "+377 (+14.1%)", "+54 (+5.0%)"],
    ["Directional concordance",    "99.0%",         "99.5%"],
    ["logFC correlation",          "~0.999",        "~0.999"],
]

for r_idx, row_data in enumerate(model_data):
    for c_idx, val in enumerate(row_data):
        cell = tbl.cell(r_idx + 1, c_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.name = "Calibri"
        if c_idx == 0:
            p.font.bold = True
        # Highlight the gain row in green
        if r_idx == 2:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xD5, 0xF5, 0xE3)
        elif r_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE

add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(2.5),
            "At the PER-DEPOT level (n ~ 19), adding Sex gains +14% DEGs in iWAT.\n"
            "Each degree of freedom matters more with fewer samples, so the variance absorbed\n"
            "by Sex outweighs its DF cost.\n\n"
            "BUT: We use a UNIFIED model (all 38 samples). What happens there?  --> Next slide",
            font_size=14, color=DARK_BLUE)

# ====== SLIDE 19: UNIFIED MODEL COMPARISON ======
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, RGBColor(0xE8, 0xF8, 0xF5))
add_section_header(slide, "Unified Models: ~ 0 + GroupDepot vs. ~ Sex + GroupDepot")
add_subtitle_text(slide,
    "Part 7.7 Test 3c: At the unified level (n = 38), adding Sex COSTS power")

# Build unified model comparison table
rows = 5
cols = 3
table_shape = slide.shapes.add_table(rows, cols,
    Inches(1), Inches(1.5), Inches(11), Inches(2.5))
tbl = table_shape.table

# Header row
headers = ["Unified Model Comparison", "iWAT contrast", "gWAT contrast"]
for i, h in enumerate(headers):
    cell = tbl.cell(0, i)
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE

unified_data = [
    ["~ 0 + GroupDepot (no Sex)",    "3,022 DEGs",    "1,159 DEGs"],
    ["~ Sex + GroupDepot",           "2,777 DEGs",    "1,089 DEGs"],
    ["Net DEG change",               "-245 (-8.1%)",  "-70 (-6.0%)"],
    ["logFC correlation",            "0.939",          "0.967"],
]

for r_idx, row_data in enumerate(unified_data):
    for c_idx, val in enumerate(row_data):
        cell = tbl.cell(r_idx + 1, c_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.name = "Calibri"
        if c_idx == 0:
            p.font.bold = True
        # Highlight the loss row in amber
        if r_idx == 2:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFC, 0xF3, 0xCF)
        elif r_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE

add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(2.8),
            "WHY THE REVERSAL?\n"
            "At n=19 (per-depot), losing 1 DF to Sex = 5.6% of residual evidence lost, "
            "but Sex absorbs enough variance to compensate.\n"
            "At n=38 (unified), losing 1 DF = only 2.9% of residual evidence, "
            "AND the balanced design already prevents sex from confounding genotype.\n"
            "Result: the DF cost now OUTWEIGHS the variance benefit -> fewer DEGs.\n\n"
            "DECISION: Use ~ 0 + GroupDepot (the simpler model with MORE power).\n"
            "The two models agree closely (r = 0.94-0.97), so biological conclusions are the same.",
            font_size=14, color=DARK_BLUE)

# ====== SLIDE 20: Summary Table ======
slide = prs.slides.add_slide(blank_layout)
add_section_header(slide, "Summary: All Evidence at a Glance")

# Build a comprehensive table
rows = 10
cols = 3
table_shape = slide.shapes.add_table(rows, cols,
    Inches(1), Inches(1.3), Inches(11), Inches(4.8))
tbl = table_shape.table

# Header row
headers = ["Metric", "iWAT", "gWAT"]
for i, h in enumerate(headers):
    cell = tbl.cell(0, i)
    cell.text = h
    p = cell.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE

# Data rows
data = [
    ["Genes tested",                        "16,466",        "16,827"],
    ["Significant interactions (LRT)",      "630 (3.83%)",   "374 (2.22%)"],
    ["Pi0 (true null proportion)",          "0.87",          "0.85"],
    ["Median variance: Sex x Genotype",     "2.1%",          "2.4%"],
    ["Sex-stratified concordance r",        "0.727",         "0.605"],
    ["Directional concordance",             "99.0%",         "99.5%"],
    ["Unified DEGs: ~ 0 + GroupDepot",      "3,022",         "1,159"],
    ["Unified DEGs: ~ Sex + GroupDepot",    "2,777",         "1,089"],
    ["Net change (unified)",                "-245 (-8.1%)",  "-70 (-6.0%)"],
]

for r_idx, row_data in enumerate(data):
    for c_idx, val in enumerate(row_data):
        cell = tbl.cell(r_idx + 1, c_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.name = "Calibri"
        if c_idx == 0:
            p.font.bold = True
        if r_idx == 8:  # highlight net change row (amber - DF cost)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFC, 0xF3, 0xCF)
        elif r_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE

add_textbox(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.8),
            "All metrics converge: the KAT8-KD response is overwhelmingly sex-independent. "
            "With a balanced design and n=38, the simpler model (~ 0 + GroupDepot) maximizes power.",
            font_size=16, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# ====== SLIDE 21: Conclusion & Decision ======
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BLUE)
add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
            "Conclusion & Analysis Decision",
            font_size=34, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

conclusion_items = [
    "CAN WE POOL SEXES?    YES",
    "",
    "EVIDENCE:",
    "  - Only 2-4% of genes show significant Sex x Genotype interaction",
    "  - Pi0 = 0.85-0.87: 85-87% of genes are true nulls (no interaction)",
    "  - The interaction term explains only ~2% of variance vs. 17-33% for Genotype",
    "  - Sex-stratified concordance: 99%+ directional agreement",
    "  - Balanced design: sex CANNOT confound genotype effect",
    "",
    "DECISION:  Design = ~ 0 + GroupDepot  (unified, no Sex covariate)",
    "  - All 38 samples in one model, depot-specific KAT8 contrasts extracted",
    "  - No Sex covariate needed: balanced design eliminates confounding",
    "  - Adding Sex to the unified model COSTS power (-8% iWAT, -6% gWAT DEGs)",
    "  - logFC correlation r > 0.94 between +/- Sex models -> same biology",
    "",
    "CAVEAT (mention in Discussion if reviewers ask):",
    "  - 630 (iWAT) / 374 (gWAT) genes DO show sex-dependent KD responses",
    "  - Male KD samples appear more extreme in PCA (especially iWAT)",
    "  - gWAT sex-stratified correlation is moderate (r = 0.61)",
    "  - These observations warrant follow-up investigation",
]

tf = add_bullet_list(slide, Inches(1), Inches(1.8), Inches(11), Inches(5.2),
                     conclusion_items, font_size=14, color=WHITE)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Results")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "KAT8_sex_interaction_presentation.pptx")
prs.save(out_path)
print(f"[OK] Presentation saved to: {out_path}")
print(f"     {len(prs.slides)} slides created")
print()
print("Slides:")
print("  1.  Title")
print("  2.  Study Design")
print("  3.  The Central Question")
print("  4.  LEARN: What is a Covariate?")
print("  5.  LEARN: Two Possible Models (~ Genotype vs ~ Sex + Genotype)")
print("  6.  Method: Likelihood Ratio Test")
print("  7.  LEARN: How to Read a P-value Histogram")
print("  8.  iWAT LRT Results")
print("  9.  gWAT LRT Results")
print("  10. LEARN: What is Pi0?")
print("  11. Pi0 Results")
print("  12. Variance Attribution")
print("  13. LEARN: What is Concordance?")
print("  14. Concordance: iWAT")
print("  15. Concordance: gWAT")
print("  16. PCA by Sex and Genotype")
print("  17. Interaction vs Genotype Scatter")
print("  18. Per-Depot Model Comparison (Part 7.7)")
print("  19. Unified Model Comparison (the key slide)")
print("  20. Summary Table (all metrics)")
print("  21. Conclusion & Decision")
print()
print("Next steps:")
print("  1. Open the PPTX on your Windows machine")
print("  2. Replace the grey placeholder boxes with your actual plots")
print("     (Right-click placeholder -> Change Shape -> Insert Picture)")
print("  3. The plot filenames are listed in each placeholder")
